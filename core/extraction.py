"""Extraction du texte d'un document (txt, pdf, pptx, image) pour alimenter l'IA."""

from pathlib import Path

import pandas as pd
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from core.gemini_client import lire_image, lire_pdf
from core.prompts import PROMPT_OCR_IMAGE, PROMPT_OCR_PDF_SCANNE

# En dessous de ce nombre de caractères, on considère qu'un PDF est un scan
# (pas de couche texte) et on le fait lire directement par Gemini.
SEUIL_TEXTE_PDF_INSUFFISANT = 50

# Gemini n'accepte les fichiers envoyés "en ligne" dans la requête (comme le fait
# cette appli, sans passer par son API de fichiers séparée) que jusqu'à ~20 Mo au
# total, encodage inclus - au-delà, la lecture échoue avec une erreur technique peu
# claire pour l'utilisateur. Streamlit, lui, autorise jusqu'à 200 Mo par fichier :
# sans cette vérification, un gros scan (fréquent pour des pages scannées en haute
# résolution) est accepté à l'upload mais échoue silencieusement à la lecture.
# Marge prise sous la limite réelle pour l'encodage et le prompt qui l'accompagne.
TAILLE_MAX_OCR_OCTETS = 15 * 1024 * 1024  # 15 Mo

MIME_PAR_EXTENSION = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
}


def type_fichier_depuis_nom(nom_fichier: str) -> str:
    ext = Path(nom_fichier).suffix.lower().lstrip(".")
    if ext == "pdf":
        return "pdf"
    if ext == "pptx":
        return "pptx"
    if ext == "docx":
        return "docx"
    if ext in ("xlsx", "csv"):
        return "tableur"
    if ext in MIME_PAR_EXTENSION:
        return "image"
    return "texte"


def extraire_texte(chemin_stocke: str, nom_original: str, api_key: str) -> str:
    """Retourne le texte extrait d'un document, quel que soit son format."""
    type_fichier = type_fichier_depuis_nom(nom_original)
    chemin = Path(chemin_stocke)

    if type_fichier == "texte":
        return chemin.read_text(encoding="utf-8", errors="replace")

    if type_fichier == "pptx":
        return _extraire_texte_pptx(chemin)

    if type_fichier == "docx":
        return _extraire_texte_docx(chemin)

    if type_fichier == "tableur":
        ext = Path(nom_original).suffix.lower().lstrip(".")
        return _extraire_texte_tableur(chemin, ext)

    if type_fichier == "image":
        ext = Path(nom_original).suffix.lower().lstrip(".")
        mime_type = MIME_PAR_EXTENSION.get(ext, "image/png")
        image_bytes = chemin.read_bytes()
        _verifier_taille_ocr(image_bytes, nom_original)
        return lire_image(image_bytes, mime_type, PROMPT_OCR_IMAGE, api_key)

    if type_fichier == "pdf":
        pdf_bytes = chemin.read_bytes()
        texte = _extraire_texte_pdf_local(chemin)
        if len(texte.strip()) < SEUIL_TEXTE_PDF_INSUFFISANT:
            # PDF probablement scanné (pas de couche texte) : on demande à Gemini de le lire.
            _verifier_taille_ocr(pdf_bytes, nom_original)
            return lire_pdf(pdf_bytes, PROMPT_OCR_PDF_SCANNE, api_key)
        return texte

    raise ValueError(f"Type de fichier non supporté : {nom_original}")


def _verifier_taille_ocr(contenu: bytes, nom_original: str):
    """Bloque AVANT d'appeler Gemini si le fichier dépasse la limite technique -
    sinon l'échec arrive côté Gemini avec un message d'erreur technique difficile
    à comprendre (ex: scan de plusieurs pages en haute résolution, très courant
    pour un document composé d'images scannées)."""
    if len(contenu) <= TAILLE_MAX_OCR_OCTETS:
        return
    taille_mo = len(contenu) / (1024 * 1024)
    limite_mo = TAILLE_MAX_OCR_OCTETS / (1024 * 1024)
    raise ValueError(
        f"fichier trop volumineux pour être lu par l'IA ({taille_mo:.1f} Mo, limite "
        f"~{limite_mo:.0f} Mo pour une image ou un PDF scanné). Essaie de le scanner "
        "en plus basse résolution, de le compresser, ou de diviser le document en "
        "plusieurs fichiers plus petits."
    )


def _extraire_texte_pdf_local(chemin: Path) -> str:
    try:
        reader = PdfReader(str(chemin))
        pages_texte = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages_texte)
    except Exception:
        return ""


def _extraire_texte_pptx(chemin: Path) -> str:
    """Lit un PowerPoint : texte de chaque diapositive (titres, contenu, tableaux)
    et notes du présentateur, directement depuis le fichier (aucun appel à l'IA
    nécessaire, c'est du texte natif dans le fichier)."""
    presentation = Presentation(str(chemin))
    morceaux = []

    for i, diapo in enumerate(presentation.slides, start=1):
        textes_diapo = []
        for forme in diapo.shapes:
            if forme.has_text_frame and forme.text_frame.text.strip():
                textes_diapo.append(forme.text_frame.text)
            if forme.has_table:
                for ligne in forme.table.rows:
                    textes_diapo.append(" | ".join(cellule.text for cellule in ligne.cells))

        if diapo.has_notes_slide:
            notes = diapo.notes_slide.notes_text_frame.text
            if notes.strip():
                textes_diapo.append(f"[Notes] {notes}")

        if textes_diapo:
            morceaux.append(f"--- Diapositive {i} ---\n" + "\n".join(textes_diapo))

    return "\n\n".join(morceaux)


def _extraire_texte_docx(chemin: Path) -> str:
    """Lit un document Word : paragraphes et tableaux, directement depuis le
    fichier (texte natif, aucun appel à l'IA nécessaire)."""
    document = Document(str(chemin))
    morceaux = [p.text for p in document.paragraphs if p.text.strip()]

    for tableau in document.tables:
        for ligne in tableau.rows:
            morceaux.append(" | ".join(cellule.text for cellule in ligne.cells))

    return "\n".join(morceaux)


def _extraire_texte_tableur(chemin: Path, ext: str) -> str:
    """Lit un fichier Excel (.xlsx) ou CSV : transforme chaque feuille en texte
    lisible par l'IA (pas besoin de vision, ce sont des données structurées)."""
    if ext == "csv":
        df = pd.read_csv(chemin)
        return df.to_string(index=False)

    feuilles = pd.read_excel(chemin, sheet_name=None)
    morceaux = []
    for nom_feuille, df in feuilles.items():
        morceaux.append(f"--- Feuille : {nom_feuille} ---\n{df.to_string(index=False)}")
    return "\n\n".join(morceaux)
